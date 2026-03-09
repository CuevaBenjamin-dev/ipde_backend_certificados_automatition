from PIL.Image import item
from fastapi import FastAPI, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from uuid import uuid4
from urllib.parse import quote
from pydantic import BaseModel
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm
from io import BytesIO
import unicodedata
import re
import os
import json
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
from copy import deepcopy
from typing import List, Tuple
from pptx.util import Pt
import qrcode
import tempfile
import subprocess

EXPORTS: dict[str, dict[str, tuple[bytes, str]]] = {}

def get_public_base_url(request: Request) -> str:
    """
    Construye la URL pública correcta usando los headers del proxy de Railway.
    Así evitamos que se generen links http:// cuando el frontend está en https://
    """
    forwarded_proto = request.headers.get("x-forwarded-proto", "")
    forwarded_host = request.headers.get("x-forwarded-host", "")
    host = request.headers.get("host", "")

    proto = (forwarded_proto.split(",")[0].strip() if forwarded_proto else request.url.scheme)
    final_host = (forwarded_host.split(",")[0].strip() if forwarded_host else host)

    # fallback seguro
    if not final_host:
        return str(request.base_url).rstrip("/").replace("http://", "https://", 1)

    return f"{proto}://{final_host}"


# -------------------------------------------------
# CONFIGURACIÓN GENERAL
# -------------------------------------------------

load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY)

app = FastAPI(title="Certificados API", version="2.3.0")

# Modelos disponibles (carpetas dentro de app/templates)
# IMPORTANTE: Deben coincidir con los valores enviados desde el frontend:
# - INSTITUTO
# - UNIVERSIDAD_2QRS
# - UNIVERSIDAD_AZUL
# - COLEGIO_ABOGADOS_CALLAO
MODELO_FOLDER_MAP = {
    "INSTITUTO": "instituto",
    "UNIVERSIDAD_2QRS": "universidad_2qrs",
    "UNIVERSIDAD_AZUL": "universidad_azul",
    "COLEGIO_ABOGADOS_CALLAO": "colegio_de_abogados_del_callao", 
    "COLEGIO_DE_PROFESORES_DEL_PERU": "colegio_de_profesores_del_peru",
}

# Modelos que usan formato de fecha largo (dd de Mes del yyyy)
MODELOS_FECHA_LARGA = {
    "INSTITUTO",
    "UNIVERSIDAD_AZUL",
    "COLEGIO_ABOGADOS_CALLAO", 
    "COLEGIO_DE_PROFESORES_DEL_PERU",
}


# Nombre de archivo de plantilla por tipo (dentro de cada carpeta de modelo)
TEMPLATE_FILENAME_MAP = {
    "DIPLOMADO": "plantilla_diplomado.pptx",
    "PROGRAMA DE ESPECIALIZACIÓN": "plantilla_programa_de_especializacion.pptx",
    "CURSO": "plantilla_curso.pptx",
    "CURSO_DE_CAPACITACION": "plantilla_curso_de_capacitacion.pptx",
    "CURSO_DE_ACTUALIZACION": "plantilla_curso_de_actualizacion.pptx",
}

# Nº de módulos por tipo
MODULOS_COUNT = {
    "DIPLOMADO": 8,
    "PROGRAMA DE ESPECIALIZACIÓN": 8,
    "CURSO": 5,
    "CURSO_DE_CAPACITACION": 5,
    "CURSO_DE_ACTUALIZACION": 5,
}

# Cache en memoria: (tipo, tema) -> módulos
MODULOS_CACHE: dict[Tuple[str, str], List[str]] = {}


# -------------------------------------------------
# UTILIDADES
# -------------------------------------------------

def safe_filename(text: str) -> str:
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^a-zA-Z0-9_-]", "_", text)
    return text


def format_date_ddmmyyyy(date_str: str) -> str:
    try:
        return datetime.strptime(date_str, "%Y-%m-%d").strftime("%d/%m/%Y")
    except ValueError:
        return date_str


def format_date_long_es(date_str: str) -> str:
    try:
        date_obj = datetime.strptime(date_str, "%Y-%m-%d")
        meses = [
            "Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
            "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"
        ]

        # ✅ día con cero a la izquierda
        dia = f"{date_obj.day:02d}"

        return f"{dia} de {meses[date_obj.month - 1]} del {date_obj.year}"
    except ValueError:
        return date_str

def format_date_range_long_es(fecha_inicio: str, fecha_fin: str) -> tuple[str, str]:
    """
    Devuelve (fecha_inicio_formateada, fecha_fin_formateada)
    aplicando la regla:
    - Si ambos años son iguales → el año solo se muestra en la fecha final
    - Si son distintos → cada fecha muestra su año
    """
    try:
        inicio = datetime.strptime(fecha_inicio, "%Y-%m-%d")
        fin = datetime.strptime(fecha_fin, "%Y-%m-%d")

        meses = [
            "Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
            "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"
        ]

        dia_inicio = f"{inicio.day:02d}"
        dia_fin = f"{fin.day:02d}"

        if inicio.year == fin.year:
            fecha_inicio_str = f"{dia_inicio} de {meses[inicio.month - 1]}"
            fecha_fin_str = f"{dia_fin} de {meses[fin.month - 1]} del {fin.year}"
        else:
            fecha_inicio_str = f"{dia_inicio} de {meses[inicio.month - 1]} del {inicio.year}"
            fecha_fin_str = f"{dia_fin} de {meses[fin.month - 1]} del {fin.year}"

        return fecha_inicio_str, fecha_fin_str

    except ValueError:
        # fallback seguro
        return format_date_long_es(fecha_inicio), format_date_long_es(fecha_fin)


def format_two_digits_number(value: int) -> str:
    """
    Formatea números enteros a dos dígitos.
    Ej: 3 -> 03, 12 -> 12
    """
    try:
        return f"{int(value):02d}"
    except (ValueError, TypeError):
        return str(value)


def format_two_digits_float(value: float) -> str:
    """
    Formatea números decimales manteniendo decimales,
    pero con parte entera a dos dígitos.
    Ej: 3 -> 03
        3.5 -> 03.5
        12 -> 12
        12.25 -> 12.25
    """
    try:
        entero = int(value)
        decimal = value - entero

        if decimal == 0:
            return f"{entero:02d}"

        # Eliminar ceros innecesarios en decimales
        decimal_str = str(round(decimal, 2)).lstrip("0")
        return f"{entero:02d}{decimal_str}"
    except (ValueError, TypeError):
        return str(value)


def calcular_horas_por_modulo(total_horas: int, cantidad_modulos: int) -> str:
    """
    Calcula horas por módulo:
    total_horas / cantidad_modulos
    Devuelve string sin decimales si es entero, o con 2 decimales si no.
    """
    if cantidad_modulos <= 0:
        return "0"

    valor = total_horas / cantidad_modulos

    # Si es entero, no mostrar decimales
    if valor.is_integer():
        return str(int(valor))

    # Si no, mostrar hasta 2 decimales
    return f"{valor:.2f}"


def nombre_completo_capitalizado(nombres: str, apellidos: str) -> str:
    texto = f"{nombres} {apellidos}".strip().lower()
    return " ".join(p.capitalize() for p in texto.split())


def modelo_con_mayuscula_inicial(tipo_modelo: str) -> str:
    texto = tipo_modelo.strip().lower().split()
    if not texto:
        return ""
    texto[0] = texto[0].capitalize()
    return " ".join(texto)


def resolve_template_path(modelo_certificado: str, tipo_modelo: str) -> str:
    modelo_key = (modelo_certificado or "").upper().strip()
    tipo_key = (tipo_modelo or "").upper().strip()

    if modelo_key not in MODELO_FOLDER_MAP:
        raise HTTPException(
            status_code=400,
            detail=f"Modelo de certificado no soportado: {modelo_key}"
        )

    if tipo_key not in TEMPLATE_FILENAME_MAP:
        raise HTTPException(
            status_code=400,
            detail=f"Tipo de certificado no soportado: {tipo_key}"
        )

    folder = MODELO_FOLDER_MAP[modelo_key]
    filename = TEMPLATE_FILENAME_MAP[tipo_key]

    return os.path.join("app", "templates", folder, filename)

##agregué
def normalize_text_for_url(text: str) -> str:
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^a-zA-Z0-9]", "", text)
    return text.lower()


def build_qr_url(nombres: str, apellidos: str, tema: str) -> str:
    primer_nombre = nombres.strip().split()[0]
    primer_apellido = apellidos.strip().split()[0]

    base = (
        normalize_text_for_url(primer_nombre)
        + normalize_text_for_url(primer_apellido)
        + normalize_text_for_url(tema)
    )

    return f"https://especializacionvirtual.com/certificados/{base}.pdf"


def generate_qr_image(url: str) -> BytesIO:
    qr = qrcode.QRCode(
        version=4,
        error_correction=qrcode.constants.ERROR_CORRECT_Q,
        box_size=10,
        border=1,
    )
    qr.add_data(url)
    qr.make(fit=True)

    img = qr.make_image(fill_color="black", back_color="white")

    buffer = BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer
##fin agregué

def insert_qr_at_placeholder(prs: Presentation, qr_stream: BytesIO):
    QR_SIZE = Cm(2.78)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            if "{{QR_CODE}}" not in shape.text_frame.text:
                continue

            left = shape.left
            top = shape.top

            shape.text_frame.clear()

            slide.shapes.add_picture(
                qr_stream,
                left=left,
                top=top,
                width=QR_SIZE,
                height=QR_SIZE,
            )

            return  # solo un QR




                    

##DISTRIBUIR HORAS POR MÓDULO
def distribuir_horas_por_modulo(total_horas: int, cantidad_modulos: int) -> List[int]:
    """
    Distribuye las horas de forma progresiva y variada,
    asegurando que la suma final sea EXACTA.
    """
    if cantidad_modulos <= 0:
        return []

    pesos = []
    incremento = 0.1
    for i in range(cantidad_modulos):
        peso = 1 + (i * incremento)
        pesos.append(peso)
    suma_pesos = sum(pesos)
    
    horas = [
        int((peso / suma_pesos) * total_horas)
        for peso in pesos
    ]

    # Ajuste para asegurar suma exacta
    diferencia = total_horas - sum(horas)

    # Repartir diferencia empezando por el último módulo
    i = cantidad_modulos - 1
    while diferencia > 0:
        horas[i] += 1
        diferencia -= 1
        i -= 1
        if i < 0:
            i = cantidad_modulos - 1

    return horas


## otro cambio
def convert_pptx_to_pdf_bytes(pptx_bytes: bytes) -> bytes:
    """
    Convierte PPTX a PDF usando LibreOffice (soffice) en modo headless.
    Requiere que 'soffice' exista en el sistema.
    """
    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = os.path.join(tmp, "input.pptx")
        pdf_path = os.path.join(tmp, "input.pdf")

        with open(pptx_path, "wb") as f:
            f.write(pptx_bytes)

        # LibreOffice crea el PDF con el mismo nombre base en el outdir
        cmd = [
            "soffice",
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--convert-to", "pdf",
            "--outdir", tmp,
            pptx_path,
        ]

        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if result.returncode != 0 or not os.path.exists(pdf_path):
            raise RuntimeError(
                "Error convirtiendo PPTX a PDF.\n"
                f"STDOUT: {result.stdout.decode(errors='ignore')}\n"
                f"STDERR: {result.stderr.decode(errors='ignore')}"
            )

        with open(pdf_path, "rb") as f:
            return f.read()




# -------------------------------------------------
# OPENAI – GENERACIÓN DE MÓDULOS DINÁMICA
# -------------------------------------------------

def build_prompt(tema: str, count: int) -> str:
    return f"""
Devuelve EXCLUSIVAMENTE un JSON válido con este formato exacto:

{{
  "modulos": [
    {', '.join(['"string"' for _ in range(count)])}
  ]
}}

REGLAS OBLIGATORIAS:
- Exactamente {count} módulos
- SOLO títulos (NO descripciones)
- Máximo 12 palabras por módulo
- En español
- NO numeración
- NO texto fuera del JSON
- NO explicaciones

Certificado: {tema}
""".strip()


def obtener_modulos_por_tema(tipo: str, tema: str) -> list[str]:
    cache_key = (tipo, tema)
    if cache_key in MODULOS_CACHE:
        return MODULOS_CACHE[cache_key]

    if tipo not in MODULOS_COUNT:
        raise ValueError("Tipo no soportado para módulos")

    count = MODULOS_COUNT[tipo]

    try:
        response = client.responses.create(
            model="gpt-5-mini",
            input=[
                {"role": "system", "content": "Eres un asistente académico extremadamente estricto."},
                {"role": "user", "content": build_prompt(tema, count)}
            ]
        )

        data = json.loads(response.output_text)
        modulos = data.get("modulos", [])

        if not isinstance(modulos, list) or len(modulos) != count:
            raise ValueError("Cantidad de módulos inválida")

        modulos = [str(m).upper().strip() for m in modulos]
        MODULOS_CACHE[cache_key] = modulos
        return modulos

    except Exception:
        return [f"MÓDULO {i+1}" for i in range(count)]
    
# -------------------------------------------------
# DIPLOMADO - AZUL - AJUSTE DE TABLA EN CERTIFICADO DE ESTUDIOS (PARTE DE ATRÁS)
# -------------------------------------------------
    
def ajustar_tabla_certificado_estudios_ua(prs: Presentation, tema: str):
    """
    Ajusta PH_TABLA en el slide donde existen PH_TEMA y PH_TABLA,
    moviendo la tabla hacia abajo según las líneas que ocupa:
    'DIPLOMADO: <tema>'
    """
    # márgenes de seguridad (en EMU)
    GAP_MIN = Cm(0.25)  # espacio mínimo entre texto y tabla (ajustable)
    MAX_EXTRA_LINES = 3  # tope lógico con tu límite front (102)

    for slide in prs.slides:
        ph_tema = find_shape_by_name(slide, "PH_TEMA")
        ph_tabla = find_shape_by_name(slide, "PH_TABLA")

        if not ph_tema or not ph_tabla:
            continue

        # 1) intentar detectar font-size real del párrafo "DIPLOMADO:"
        font_pt = 12.0
        try:
            tf = ph_tema.text_frame
            for p in tf.paragraphs:
                full = "".join(run.text for run in p.runs)
                if "DIPLOMADO" in full:
                    # agarra la primera run con tamaño definido
                    for run in p.runs:
                        if run.font and run.font.size:
                            font_pt = run.font.size.pt
                            raise StopIteration
        except StopIteration:
            pass
        except Exception:
            pass

        # 2) calcular cuántas líneas ocuparía la línea DIPLOMADO:
        texto_linea = f"DIPLOMADO: {tema}".strip()

        max_chars = estimate_chars_per_line(ph_tema.width, font_pt)

        # OJO: en tu template hay también "PARTICIPANTE: ...", pero nosotros
        # solo usamos esta línea para decidir el empuje hacia abajo.
        lines = wrap_by_words(texto_linea, max_chars)
        line_count = len(lines)

        # 3) si usa 1 línea, no movemos nada
        if line_count <= 1:
            return

        # protección: no debería pasar con el front en 102, pero por si acaso
        if line_count > MAX_EXTRA_LINES:
            line_count = MAX_EXTRA_LINES

        # 4) calcular delta (alto de línea aproximado)
        line_height_emu = int((font_pt * 0.65) * EMU_PER_PT)  # 0.65 es típico
        extra_lines = line_count - 1
        delta = extra_lines * line_height_emu

        PUSH_EXTRA = Cm(0.00)
        if line_count >= 3:
            PUSH_EXTRA = Cm(0.25)

        ph_tabla.top = ph_tabla.top + delta + int(GAP_MIN) + int(PUSH_EXTRA)

        return  # ajustamos una sola vez

# -------------------------------------------------
# MISMA LÓGICA DE ARRIBA PERO PARA PROGRAMA DE ESPECIALIZACIÓN (CASO PARTICULAR DE UNIVERSIDAD AZUL)
# -------------------------------------------------
    
def ajustar_tabla_certificado_estudios_ua_programa(prs: Presentation, tema: str):
    """
    Ajusta PH_TABLA en el slide donde existen PH_TEMA y PH_TABLA,
    moviendo la tabla hacia abajo según las líneas que ocupa:
    'PROGRAMA DE ESPECIALIZACIÓN: <tema>'
    """
    GAP_MIN = Cm(0.25)       # espacio mínimo entre texto y tabla
    MAX_EXTRA_LINES = 3      # tu tope lógico

    for slide in prs.slides:
        ph_tema = find_shape_by_name(slide, "PH_TEMA")
        ph_tabla = find_shape_by_name(slide, "PH_TABLA")

        if not ph_tema or not ph_tabla:
            continue

        # 1) detectar font-size real del párrafo "PROGRAMA..."
        font_pt = 12.0
        try:
            tf = ph_tema.text_frame
            for p in tf.paragraphs:
                full = "".join(run.text for run in p.runs)
                if "PROGRAMA" in full:  # 👈 clave para programa
                    for run in p.runs:
                        if run.font and run.font.size:
                            font_pt = run.font.size.pt
                            raise StopIteration
        except StopIteration:
            pass
        except Exception:
            pass

        # 2) calcular cuántas líneas ocuparía la línea PROGRAMA:
        texto_linea = f"PROGRAMA DE ESPECIALIZACIÓN: {tema}".strip()

        max_chars = estimate_chars_per_line(ph_tema.width, font_pt)
        lines = wrap_by_words(texto_linea, max_chars)
        line_count = len(lines)

        # 3) si usa 1 línea, no movemos nada
        if line_count <= 1:
            return

        # protección
        if line_count > MAX_EXTRA_LINES:
            line_count = MAX_EXTRA_LINES

        # 4) delta por líneas extra
        line_height_emu = int((font_pt * 0.65) * EMU_PER_PT)
        extra_lines = line_count - 1
        delta = extra_lines * line_height_emu

        # 5) empuje extra cuando ya estamos en 3 líneas (igual que diplomado)
        PUSH_EXTRA = Cm(0.00)
        if line_count >= 3:
            PUSH_EXTRA = Cm(0.25)

        ph_tabla.top = ph_tabla.top + delta + int(GAP_MIN) + int(PUSH_EXTRA)
        return
    

# -------------------------------------------------
# MISMA LÓGICA PERO PARA CURSO (UNIVERSIDAD AZUL)
# -------------------------------------------------

def ajustar_tabla_certificado_estudios_ua_curso(prs: Presentation, tema: str):
    """
    Ajusta PH_TABLA en el slide donde existen PH_TEMA y PH_TABLA,
    moviendo la tabla hacia abajo según las líneas que ocupa:
    'CURSO: <tema>'
    """
    GAP_MIN = Cm(0.25)       # espacio mínimo entre PH_TEMA y PH_TABLA
    MAX_EXTRA_LINES = 3      # tope lógico

    for slide in prs.slides:
        ph_tema = find_shape_by_name(slide, "PH_TEMA")
        ph_tabla = find_shape_by_name(slide, "PH_TABLA")

        if not ph_tema or not ph_tabla:
            continue

        # 1) detectar font-size real del párrafo que contiene "CURSO"
        font_pt = 12.0
        try:
            tf = ph_tema.text_frame
            for p in tf.paragraphs:
                full = "".join(run.text for run in p.runs)
                if "CURSO" in full:
                    for run in p.runs:
                        if run.font and run.font.size:
                            font_pt = run.font.size.pt
                            raise StopIteration
        except StopIteration:
            pass
        except Exception:
            pass

        # 2) calcular cuántas líneas ocuparía la línea CURSO:
        texto_linea = f"CURSO: {tema}".strip()

        max_chars = estimate_chars_per_line(ph_tema.width, font_pt)
        lines = wrap_by_words(texto_linea, max_chars)
        line_count = len(lines)

        # 3) si usa 1 línea, no movemos nada
        if line_count <= 1:
            return

        # protección
        if line_count > MAX_EXTRA_LINES:
            line_count = MAX_EXTRA_LINES

        # 4) delta por líneas extra
        line_height_emu = int((font_pt * 0.65) * EMU_PER_PT)
        extra_lines = line_count - 1
        delta = extra_lines * line_height_emu

        # 5) empuje extra cuando ya estamos en 3 líneas (igual que los otros)
        PUSH_EXTRA = Cm(0.00)
        if line_count >= 3:
            PUSH_EXTRA = Cm(0.25)

        ph_tabla.top = ph_tabla.top + delta + int(GAP_MIN) + int(PUSH_EXTRA)
        return  # ajustamos una sola vez


# -------------------------------------------------
# MISMA LÓGICA PERO PARA CURSO DE CAPACITACIÓN (UNIVERSIDAD AZUL)
# -------------------------------------------------

def ajustar_tabla_certificado_estudios_ua_curso_capacitacion(prs: Presentation, tema: str):
    """
    Ajusta PH_TABLA en el slide donde existen PH_TEMA y PH_TABLA,
    moviendo la tabla hacia abajo según las líneas que ocupa:
    'CURSO DE CAPACITACIÓN: <tema>'
    """
    GAP_MIN = Cm(0.25)       # espacio mínimo entre PH_TEMA y PH_TABLA
    MAX_EXTRA_LINES = 3      # tope lógico

    for slide in prs.slides:
        ph_tema = find_shape_by_name(slide, "PH_TEMA")
        ph_tabla = find_shape_by_name(slide, "PH_TABLA")

        if not ph_tema or not ph_tabla:
            continue

        # 1) detectar font-size real del párrafo que contiene "CAPACITACIÓN"
        font_pt = 12.0
        try:
            tf = ph_tema.text_frame
            for p in tf.paragraphs:
                full = "".join(run.text for run in p.runs)
                # buscamos algo distintivo de este tipo
                if "CAPACIT" in full or "CURSO" in full:
                    for run in p.runs:
                        if run.font and run.font.size:
                            font_pt = run.font.size.pt
                            raise StopIteration
        except StopIteration:
            pass
        except Exception:
            pass

        # 2) calcular cuántas líneas ocuparía la línea CURSO DE CAPACITACIÓN:
        texto_linea = f"CURSO DE CAPACITACIÓN: {tema}".strip()

        max_chars = estimate_chars_per_line(ph_tema.width, font_pt)
        lines = wrap_by_words(texto_linea, max_chars)
        line_count = len(lines)

        # 3) si usa 1 línea, no movemos nada
        if line_count <= 1:
            return

        # protección
        if line_count > MAX_EXTRA_LINES:
            line_count = MAX_EXTRA_LINES

        # 4) delta por líneas extra
        line_height_emu = int((font_pt * 0.65) * EMU_PER_PT)
        extra_lines = line_count - 1
        delta = extra_lines * line_height_emu

        # 5) empuje extra cuando ya estamos en 3 líneas
        PUSH_EXTRA = Cm(0.00)
        if line_count >= 3:
            PUSH_EXTRA = Cm(0.25)

        ph_tabla.top = ph_tabla.top + delta + int(GAP_MIN) + int(PUSH_EXTRA)
        return  # ajustamos una sola vez


# -------------------------------------------------
# MISMA LÓGICA PERO PARA CURSO DE ACTUALIZACIÓN (UNIVERSIDAD AZUL)
# -------------------------------------------------

def ajustar_tabla_certificado_estudios_ua_curso_actualizacion(prs: Presentation, tema: str):
    """
    Ajusta PH_TABLA en el slide donde existen PH_TEMA y PH_TABLA,
    moviendo la tabla hacia abajo según las líneas que ocupa:
    'CURSO DE ACTUALIZACIÓN: <tema>'
    """
    GAP_MIN = Cm(0.25)       # espacio mínimo entre PH_TEMA y PH_TABLA
    MAX_EXTRA_LINES = 3      # tope lógico

    for slide in prs.slides:
        ph_tema = find_shape_by_name(slide, "PH_TEMA")
        ph_tabla = find_shape_by_name(slide, "PH_TABLA")

        if not ph_tema or not ph_tabla:
            continue

        # 1) detectar font-size real del párrafo que contiene "ACTUALIZACIÓN"
        font_pt = 12.0
        try:
            tf = ph_tema.text_frame
            for p in tf.paragraphs:
                full = "".join(run.text for run in p.runs)
                # distintivo: ACTUALIZACIÓN
                if "ACTUALIZ" in full or "CURSO" in full:
                    for run in p.runs:
                        if run.font and run.font.size:
                            font_pt = run.font.size.pt
                            raise StopIteration
        except StopIteration:
            pass
        except Exception:
            pass

        # 2) calcular cuántas líneas ocuparía la línea CURSO DE ACTUALIZACIÓN:
        texto_linea = f"CURSO DE ACTUALIZACIÓN: {tema}".strip()

        max_chars = estimate_chars_per_line(ph_tema.width, font_pt)
        lines = wrap_by_words(texto_linea, max_chars)
        line_count = len(lines)

        # 3) si usa 1 línea, no movemos nada
        if line_count <= 1:
            return

        # protección
        if line_count > MAX_EXTRA_LINES:
            line_count = MAX_EXTRA_LINES

        # 4) delta por líneas extra
        line_height_emu = int((font_pt * 0.65) * EMU_PER_PT)
        extra_lines = line_count - 1
        delta = extra_lines * line_height_emu

        # 5) empuje extra cuando ya estamos en 3 líneas
        PUSH_EXTRA = Cm(0.00)
        if line_count >= 3:
            PUSH_EXTRA = Cm(0.25)

        ph_tabla.top = ph_tabla.top + delta + int(GAP_MIN) + int(PUSH_EXTRA)
        return  # ajustamos una sola vez
    

# -------------------------------------------------
# FUNCIÓN GENÉRICA: AJUSTAR TABLA SEGÚN LÍNEAS DEL TIPO + TEMA
# ESTO SIRVE PARA MODELOS 2QRS
# (sirve para cualquier modelo mientras existan PH_TEMA y PH_TABLA)
# -------------------------------------------------

def ajustar_tabla_certificado_estudios_generico(
    prs: Presentation,
    label: str,
    tema: str,
    *,
    shape_tema_name: str = "PH_TEMA",
    shape_tabla_name: str = "PH_TABLA",
    gap_min: Cm = Cm(0.25),
    max_extra_lines: int = 3,
):
    """
    Mueve hacia abajo la tabla (PH_TABLA) si el texto:
        '<LABEL>: <TEMA>'
    ocupa más de 1 línea dentro de PH_TEMA.

    - label: Ej: 'DIPLOMADO', 'PROGRAMA DE ESPECIALIZACIÓN', 'CURSO', etc.
    - tema: el tema ya en mayúsculas (o como lo uses)
    """
    label = (label or "").strip().upper()
    tema = (tema or "").strip()

    # texto que realmente evaluamos (simulación del wrap)
    texto_linea = f"{label}: {tema}".strip()

    for slide in prs.slides:
        ph_tema = find_shape_by_name(slide, shape_tema_name)
        ph_tabla = find_shape_by_name(slide, shape_tabla_name)

        if not ph_tema or not ph_tabla:
            continue

        # 1) intentar detectar font-size real del párrafo que contiene el label
        font_pt = 12.0
        try:
            tf = ph_tema.text_frame
            for p in tf.paragraphs:
                full = "".join(run.text for run in p.runs).upper()
                if label and label in full:
                    for run in p.runs:
                        if run.font and run.font.size:
                            font_pt = run.font.size.pt
                            raise StopIteration
        except StopIteration:
            pass
        except Exception:
            pass

        # 2) calcular cuántas líneas ocuparía esa línea
        max_chars = estimate_chars_per_line(ph_tema.width, font_pt)
        lines = wrap_by_words(texto_linea, max_chars)
        line_count = len(lines)

        # 3) si usa 1 línea, no movemos nada
        if line_count <= 1:
            return

        # protección
        if line_count > max_extra_lines:
            line_count = max_extra_lines

        # 4) delta por líneas extra
        line_height_emu = int((font_pt * 0.65) * EMU_PER_PT)
        extra_lines = line_count - 1
        delta = extra_lines * line_height_emu

        # empuje extra cuando ya estamos en 3 líneas
        push_extra = Cm(0.00)
        if line_count >= 3:
            push_extra = Cm(0.25)

        ph_tabla.top = ph_tabla.top + delta + int(gap_min) + int(push_extra)
        return  # ajustamos una sola vez (primer slide que lo encuentre)


# -------------------------------------------------
# REGLAS DE AJUSTE PARA CERTIFICADO DE ESTUDIOS (PARTE DE ATRÁS)
# -------------------------------------------------

AJUSTES_TABLA_POR_MODELO_Y_TIPO = {
    "UNIVERSIDAD_AZUL": {
        "DIPLOMADO": "DIPLOMADO",
        "PROGRAMA DE ESPECIALIZACIÓN": "PROGRAMA DE ESPECIALIZACIÓN",
        "CURSO": "CURSO",
        "CURSO_DE_CAPACITACION": "CURSO DE CAPACITACIÓN",
        "CURSO_DE_ACTUALIZACION": "CURSO DE ACTUALIZACIÓN",
    },
    "UNIVERSIDAD_2QRS": {
        "DIPLOMADO": "DIPLOMADO",
        "PROGRAMA DE ESPECIALIZACIÓN": "PROGRAMA DE ESPECIALIZACIÓN",
        "CURSO": "CURSO",
        "CURSO_DE_CAPACITACION": "CURSO DE CAPACITACIÓN",
        "CURSO_DE_ACTUALIZACION": "CURSO DE ACTUALIZACIÓN",
    },
        "INSTITUTO": {
        "DIPLOMADO": "DIPLOMADO",
        "PROGRAMA DE ESPECIALIZACIÓN": "PROGRAMA DE ESPECIALIZACIÓN",
        "CURSO": "CURSO",
        "CURSO_DE_CAPACITACION": "CURSO DE CAPACITACIÓN",
        "CURSO_DE_ACTUALIZACION": "CURSO DE ACTUALIZACIÓN",
    },
}


# -------------------------------------------------
# CORS
# -------------------------------------------------

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://ipde-frontend-certificado-automatit.vercel.app",
        "http://localhost:4200",
    ],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition"],
)

@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    return JSONResponse(
        status_code=500,
        content={"detail": str(exc)}
    )


# -------------------------------------------------
# DTO
# -------------------------------------------------

## otro cambio
class DiplomaRequest(BaseModel):
    modeloCertificado: str
    tipoModelo: str
    nombres: str
    apellidos: str
    temaDiplomado: str
    fechaInicio: str
    fechaFin: str
    horasAcademicas: int
    creditosAcademicos: int
    folioNumero: str
    fechaEmision: str
    codigoEstudiante: str = ""
    qrSlug: str = ""

class BatchRequest(BaseModel):
    items: List[DiplomaRequest]


# -------------------------------------------------
# REEMPLAZO DE TEXTO (PRESERVA FUENTES)
# -------------------------------------------------

def replace_in_text_frame_preserve_font(text_frame, mapping):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            for k, v in mapping.items():
                if k in run.text:
                    run.text = run.text.replace(k, v)


def replace_in_shape(shape, mapping):
    if shape.has_text_frame:
        replace_in_text_frame_preserve_font(shape.text_frame, mapping)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_in_text_frame_preserve_font(cell.text_frame, mapping)

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            replace_in_shape(s, mapping)


def replace_placeholders(prs, mapping):
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape, mapping)

        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                replace_in_shape(shape, mapping)

    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            replace_in_shape(shape, mapping)

    for master in prs.slide_masters:
        for shape in master.shapes:
            replace_in_shape(shape, mapping)
            
            
EMU_PER_PT = 12700

def find_shape_by_name(slide, name: str):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    return None

def estimate_chars_per_line(shape_width_emu: int, font_size_pt: float) -> int:
    """
    Estimación razonable sin medir fuente real:
    - ancho en pt = width_emu / 12700
    - ancho promedio de carácter ~ 0.52 * font_size (depende de la fuente)
    """
    if not font_size_pt or font_size_pt <= 0:
        font_size_pt = 12.0

    width_pt = shape_width_emu / EMU_PER_PT
    avg_char_pt = font_size_pt * 0.52
    cpl = int(width_pt / avg_char_pt)

    SAFETY_FACTOR = 0.85  # prueba 0.88; si aún falla a 3 líneas, baja a 0.85
    cpl = int(cpl * SAFETY_FACTOR)

    return max(10, cpl)

def wrap_by_words(text: str, max_chars: int) -> list[str]:
    """
    Simula word-wrap:
    - arma líneas sin partir palabras (salvo palabras largas)
    - si una palabra es más larga que max_chars, se parte (para no romper todo)
    """
    words = text.split(" ")
    lines = []
    current = ""

    for w in words:
        if current == "":
            candidate = w
        else:
            candidate = current + " " + w

        if len(candidate) <= max_chars:
            current = candidate
            continue

        # si current ya tiene algo, lo cerramos
        if current:
            lines.append(current)
            current = ""

        # si la palabra sola no cabe, la partimos
        while len(w) > max_chars:
            lines.append(w[:max_chars])
            w = w[max_chars:]

        current = w

    if current:
        lines.append(current)

    return lines


# -------------------------------------------------
# MERGE PPTX
# -------------------------------------------------

def _replace_rids_in_element(el, rid_map: dict[str, str]):
    for e in el.iter():
        if not hasattr(e, "attrib"):
            continue
        for attr_key, attr_val in list(e.attrib.items()):
            if attr_val in rid_map:
                e.attrib[attr_key] = rid_map[attr_val]


def clone_slide_into(dest_prs: Presentation, src_slide):
    # ✅ USAR LAYOUT EN BLANCO (evita placeholders duplicados)
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    spTree = new_slide.shapes._spTree
    src_spTree = src_slide.shapes._spTree

    for child in list(src_spTree):
        tag = child.tag.lower()

        # ❌ NO copiar propiedades internas del layout
        if tag.endswith("nvgrpsppr") or tag.endswith("grpsppr"):
            continue

        spTree.insert_element_before(deepcopy(child), 'p:extLst')

    # 🔁 Copiar relaciones (imágenes, fondos, etc.)
    rid_map = {}
    for rId, rel in src_slide.part.rels.items():
        try:
            new_rId = new_slide.part.relate_to(
                rel._target,
                rel.reltype,
                is_external=rel.is_external
            )
            rid_map[rId] = new_rId
        except Exception:
            continue

    _replace_rids_in_element(new_slide._element, rid_map)


def merge_presentations(presentations: List[Presentation]) -> Presentation:
    if not presentations:
        raise ValueError("No hay presentaciones para unir")

    # 🔒 Usar la primera presentación como base (PRESERVA TEMA Y COLORES)
    dest = presentations[0]

    for prs in presentations[1:]:
        for slide in prs.slides:
            clone_slide_into(dest, slide)

    return dest


# -------------------------------------------------
# GENERACIÓN DE PPTX POR ITEM  ✅ RESTAURADA
# -------------------------------------------------

def generar_presentacion_por_item(item: DiplomaRequest) -> Presentation:
    tipo = item.tipoModelo.upper().strip()
    modelo_cert = item.modeloCertificado.upper().strip()

    template_path = resolve_template_path(modelo_cert, tipo)

    if not os.path.exists(template_path):
        raise HTTPException(
            status_code=500,
            detail=f"No existe la plantilla: {template_path}"
        )

    prs = Presentation(template_path)
    
    # Determinar formato de fecha según modelo de certificado
    usar_fecha_larga = modelo_cert in MODELOS_FECHA_LARGA

    if usar_fecha_larga:
        fecha_inicio, fecha_fin = format_date_range_long_es(
            item.fechaInicio,
            item.fechaFin
        )
        fecha_emision_larga = format_date_long_es(item.fechaEmision)
        fecha_emision_corta = format_date_ddmmyyyy(item.fechaEmision)
        # fecha_emision = format_date_long_es(item.fechaEmision)
    else:
        fecha_inicio = format_date_ddmmyyyy(item.fechaInicio)
        fecha_fin = format_date_ddmmyyyy(item.fechaFin)
        fecha_emision_larga = format_date_long_es(item.fechaEmision)
        fecha_emision_corta = format_date_ddmmyyyy(item.fechaEmision)
        # fecha_emision = format_date_ddmmyyyy(item.fechaEmision)
    
    nombre_posterior = nombre_completo_capitalizado(item.nombres, item.apellidos)
    modulos = obtener_modulos_por_tema(tipo, item.temaDiplomado)
    
    ##horas de módulos correctamente distribuidas
    cantidad_modulos = MODULOS_COUNT[tipo]

    horas_por_modulo_global = calcular_horas_por_modulo(
        item.horasAcademicas,
        cantidad_modulos
    )
    
    horas_por_modulo = distribuir_horas_por_modulo(
        item.horasAcademicas,
        cantidad_modulos
    )

    mapping = {
        "{{MODELO}}": tipo,
        "{{MODELO_MINUSCULA}}": modelo_con_mayuscula_inicial(tipo),

        "{{NOMBRES}}": item.nombres.upper(),
        "{{APELLIDOS}}": item.apellidos.upper(),
        "{{TEMA_DIPLOMADO}}": item.temaDiplomado.upper(),

        "{{NOMBRE_COMPLETO_MINUSCULA}}": nombre_posterior,

        "{{FECHA_INICIO}}": fecha_inicio,
        "{{FECHA_FIN}}": fecha_fin,

        "{{HORAS_ACADEMICAS}}": format_two_digits_number(item.horasAcademicas),
        "{{CREDITOS_ACADEMICOS}}": format_two_digits_float(item.creditosAcademicos),
        
        "{{HORAS_MODULO}}": horas_por_modulo_global,

        "{{FOLIO_NUMERO}}": item.folioNumero,

        "{{FECHA_EMISION_LARGA}}": fecha_emision_larga,
        "{{FECHA_EMISION_CORTA}}": fecha_emision_corta, 
        
        "{{CODE_STUDENT}}": (item.codigoEstudiante or "").strip(),
    }

    for i, m in enumerate(modulos, start=1):
        mapping[f"{{{{MODULO_{i}}}}}"] = m
    
    ##Agregar horas por módulo al mapping       
    for i, h in enumerate(horas_por_modulo, start=1):
        mapping[f"{{{{HORAS_MODULO_{i}}}}}"] = str(h)


    replace_placeholders(prs, mapping)


# -------------------------------------------------
# AJUSTE GENÉRICO (AZUL y 2QRS) - CERTIFICADO DE ESTUDIOS / PARTE DE ATRÁS
# -------------------------------------------------

    modelo_ajustes = AJUSTES_TABLA_POR_MODELO_Y_TIPO.get(modelo_cert)
    if modelo_ajustes:
        label = modelo_ajustes.get(tipo)
        if label:
            ajustar_tabla_certificado_estudios_generico(
                prs,
                label=label,
                tema=item.temaDiplomado.upper(),
            )

    return prs


# -------------------------------------------------
# ENDPOINTS
# -------------------------------------------------

@app.get("/health")
def health():
    return {"status": "ok"}

@app.post("/api/diplomas")
def generate_pptx_batch(payload: BatchRequest, request: Request):
    if not payload.items:
        raise HTTPException(status_code=400, detail="items no puede estar vacío")

    # 1) Generar presentaciones por item
    presentations: List[Presentation] = []
    for item in payload.items:
        prs = generar_presentacion_por_item(item)
        presentations.append(prs)

    # 2) Merge PPTX
    merged = merge_presentations(presentations)

    # 3) Insertar QR en el merged (SE DEJA TAL CUAL)
    for item in payload.items:
        slug = (item.qrSlug or "").strip()
        if slug:
            qr_url = f"https://especializacionvirtual.com/certificados/{slug}.pdf"
        else:
            qr_url = build_qr_url(item.nombres, item.apellidos, item.temaDiplomado)

        qr_image = generate_qr_image(qr_url)
        insert_qr_at_placeholder(merged, qr_image)

    # 4) Guardar PPTX final a bytes
    merged_buf = BytesIO()
    merged.save(merged_buf)
    merged_pptx_bytes = merged_buf.getvalue()

    # 5) Generar PDFs por cada item
    pdf_files: list[tuple[str, bytes]] = []

    for i, (prs_item, item) in enumerate(zip(presentations, payload.items), start=1):
        b = BytesIO()
        prs_item.save(b)
        pptx_bytes = b.getvalue()

        pdf_bytes = convert_pptx_to_pdf_bytes(pptx_bytes)

        slug = (item.qrSlug or "").strip()
        if not slug:
            slug = (
                normalize_text_for_url(item.nombres.split()[0])
                + normalize_text_for_url(item.apellidos.split()[0])
                + normalize_text_for_url(item.temaDiplomado)
            )

        pdf_name = safe_filename(slug) + ".pdf"
        pdf_files.append((pdf_name, pdf_bytes))

    # 6) Guardar temporalmente archivos en memoria
    export_id = uuid4().hex
    files_to_store: dict[str, tuple[bytes, str]] = {}

    pptx_name = f"CERTIFICADOS_{int(datetime.now().timestamp())}.pptx"
    files_to_store[pptx_name] = (
        merged_pptx_bytes,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    for name, content in pdf_files:
        files_to_store[name] = (content, "application/pdf")

    EXPORTS[export_id] = files_to_store

    # 7) Construir URLs públicas de descarga
    base_url = get_public_base_url(request)
    response_files = []

    for name in files_to_store.keys():
        response_files.append({
            "name": name,
            "url": f"{base_url}/api/diplomas/download/{export_id}/{quote(name)}",
        })

    # 8) Devolver JSON con los archivos a descargar
    return {
        "exportId": export_id,
        "files": response_files,
    }

@app.get("/api/diplomas/download/{export_id}/{filename:path}")
def download_generated_file(export_id: str, filename: str):
    export = EXPORTS.get(export_id)
    if not export:
        raise HTTPException(status_code=404, detail="Exportación no encontrada")

    file_data = export.get(filename)
    if not file_data:
        raise HTTPException(status_code=404, detail="Archivo no encontrado")

    content, media_type = file_data

    return StreamingResponse(
        BytesIO(content),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )